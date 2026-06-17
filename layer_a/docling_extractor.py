"""
docling_extractor.py — Docling Extractor（schema-v3.0）

轉換層職責：PDF → Docling 原始結構，最小轉換。
和 azure_cu_extractor 一樣：直接保留工具的原始輸出。

Docling export_to_dict() 原始輸出直接保留：
  texts[]        — 所有文字項目（含 bbox / charspan / label / formatting）
  tables[]       — 所有表格（含 cell-level data：bbox / row_span / col_span / header）
  pictures[]     — 所有圖片（含 bbox，裁切比 CU API 更精確）
  pages{}        — 頁面資訊（含 size）
  groups[]       — 分組結構
  key_value_items[] — 鍵值對
  body           — 文件樹根節點
  furniture      — 頁首/頁尾等非正文元素

另外計算：
  markdown  — Docling export_to_markdown()
  confidence — source=pdf_text_layer，available=False
  qc         — 資訊損失量化
  metadata   — 標準四層
"""

from __future__ import annotations
import sys, hashlib
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))


# ── 圖片裁切（使用 Docling 提供的精確 bbox）────────────────────────────────

def _materialize_pictures(
    pdf_path: Path,
    pictures: list[dict],
    output_dir: Path,
    page_sizes: dict,
) -> list[dict]:
    """用 Docling bbox 裁切並存圖（比 CU API 更精確：PictureItem bbox = 實際圖範圍）。"""
    import fitz

    img_dir = output_dir / "figures"
    img_dir.mkdir(parents=True, exist_ok=True)

    doc = fitz.open(str(pdf_path))
    result = []

    for seq, pic in enumerate(pictures, 1):
        pic_out = dict(pic)
        provs   = pic.get("prov", [])
        if not provs:
            pic_out.update({"path": None, "sha256": None, "has_image": False})
            result.append(pic_out)
            continue

        prov    = provs[0]
        page_no = prov.get("page_no", 1)
        bbox    = prov.get("bbox", {})
        if page_no - 1 >= len(doc) or not bbox:
            pic_out.update({"path": None, "sha256": None, "has_image": False})
            result.append(pic_out)
            continue

        page = doc[page_no - 1]
        l, t, r, b = bbox["l"], bbox["t"], bbox["r"], bbox["b"]

        # Docling BOTTOMLEFT → fitz TOPLEFT 座標系轉換
        # normalize 防止座標反轉（某些 PDF 工具輸出異常 bbox）
        try:
            if bbox.get("coord_origin", "BOTTOMLEFT") == "BOTTOMLEFT":
                ph = page_sizes.get(str(page_no), {}).get("height", page.rect.height)
                x0, x1 = min(l, r), max(l, r)
                y0, y1 = min(ph - t, ph - b), max(ph - t, ph - b)
            else:
                x0, x1 = min(l, r), max(l, r)
                y0, y1 = min(t, b), max(t, b)
            rect = fitz.Rect(x0, y0, x1, y1) & page.rect
        except Exception:
            pic_out.update({"path": None, "sha256": None, "has_image": False})
            result.append(pic_out)
            continue

        if rect.is_empty:
            pic_out.update({"path": None, "sha256": None, "has_image": False})
            result.append(pic_out)
            continue

        pixmap    = page.get_pixmap(clip=rect, dpi=150)
        png_bytes = pixmap.tobytes("png")
        sha256    = hashlib.sha256(png_bytes).hexdigest()
        fname     = f"{pdf_path.stem}_p{page_no}_pic{seq}.png"
        fpath     = img_dir / fname
        fpath.write_bytes(png_bytes)

        pic_out.update({
            "path":      str(fpath.relative_to(output_dir)),
            "sha256":    sha256,
            "has_image": True,
        })
        result.append(pic_out)

    doc.close()
    return result


def _materialize_pictures_from_uri(
    pictures: list[dict],
    output_dir: Path,
    stem: str,
) -> list[dict]:
    """DOCX path：從 Docling 內嵌的 base64 data URI 解碼存圖。

    Docling 對 DOCX 的 pictures[].image.uri 格式為
    'data:<mimetype>;base64,<data>'，直接解碼即可取得圖片位元組。
    不需要 fitz（PDF-only）。
    """
    import base64 as _b64
    img_dir = output_dir / "figures"
    img_dir.mkdir(parents=True, exist_ok=True)
    result = []
    for seq, pic in enumerate(pictures, 1):
        pic_out = dict(pic)
        uri = (pic.get("image") or {}).get("uri", "")
        if not uri.startswith("data:"):
            pic_out.update({"path": None, "sha256": None, "has_image": False})
            result.append(pic_out)
            continue
        try:
            header, b64data = uri.split(",", 1)
            mimetype = header.split(";")[0].replace("data:", "")
            ext = mimetype.split("/")[-1] if "/" in mimetype else "png"
            img_bytes = _b64.b64decode(b64data)
            sha256 = hashlib.sha256(img_bytes).hexdigest()
            fpath = img_dir / f"{stem}_pic{seq}.{ext}"
            fpath.write_bytes(img_bytes)
            pic_out.update({
                "path":      str(fpath.relative_to(output_dir)),
                "sha256":    sha256,
                "has_image": True,
            })
        except Exception:
            pic_out.update({"path": None, "sha256": None, "has_image": False})
        result.append(pic_out)
    return result


# ── QC 計算 ──────────────────────────────────────────────────────────────────

import re as _re
# 私用區（U+E000-F8FF）+ 替換字元（U+FFFD），不含 CJK 相容字（U+F900-FAFF）
_PUA_RE = _re.compile(r"[-�]")


def _compute_qc(raw: dict, page_count: int,
                pictures: list[dict], output_dir_set: bool) -> dict:
    texts  = raw.get("texts",  [])
    tables = raw.get("tables", [])

    # ── 1. 空白頁 & 低密度頁 ──────────────────────────────────────────────
    items_per_page_map: dict[int, int] = {}
    for key in ("texts", "tables", "pictures"):
        for item in raw.get(key, []):
            for prov in item.get("prov", []):
                pn = prov.get("page_no", 0)
                items_per_page_map[pn] = items_per_page_map.get(pn, 0) + 1

    all_page_nos = set(range(1, page_count + 1))
    # "blank" 在此指「Docling 未萃取出任何 item 的頁面」，不代表頁面真正空白
    pages_no_extracted_items = sorted(all_page_nos - set(items_per_page_map))
    total_items  = sum(items_per_page_map.values())
    avg_density  = total_items / max(page_count, 1)
    low_threshold = max(1, avg_density * 0.3)
    low_density_pages = sorted(
        pn for pn, cnt in items_per_page_map.items() if cnt < low_threshold
    )

    # ── 2. 亂碼字元率（PUA / 亂碼字元 / 總字元）─────────────────────────
    chars_total   = sum(len(t.get("text") or "") for t in texts)
    chars_garbled = sum(
        len(_PUA_RE.findall(t.get("text") or "")) for t in texts
    )
    garbled_char_rate = round(chars_garbled / max(chars_total, 1), 4)

    # ── 3. 表格空白 cell 率 ───────────────────────────────────────────────
    cells_total = 0
    cells_empty = 0
    for tbl in tables:
        for cell in tbl.get("data", {}).get("table_cells", []):
            cells_total += 1
            if not (cell.get("text") or "").strip():
                cells_empty += 1
    empty_cell_rate = round(cells_empty / max(cells_total, 1), 4)

    # ── 4. 圖片覆蓋 ───────────────────────────────────────────────────────
    figures_total        = len(pictures)
    figures_materialized = sum(1 for p in pictures if p.get("has_image"))
    figure_loss = (
        round(1 - figures_materialized / figures_total, 4)
        if output_dir_set and figures_total > 0 else 0.0
    )

    # ── 5. 綜合損失率 ─────────────────────────────────────────────────────
    text_loss  = garbled_char_rate
    page_loss  = round(len(low_density_pages) / max(page_count, 1), 4)
    table_loss = empty_cell_rate

    # 啟發式加權公式（非精確損失率）：係數依文件內容對 RAG 的重要性估計，無實驗驗證
    estimated = round(
        0.45 * text_loss +
        0.25 * page_loss +
        0.20 * table_loss +
        0.10 * figure_loss,
        4
    )
    qc_level = (
        "danger"  if estimated > 0.10 else
        "warning" if estimated > 0.03 else
        "good"
    )

    return {
        # 文字品質
        "garbled_char_rate":         garbled_char_rate,
        # 頁面密度
        "items_per_page":            round(avg_density, 1),
        "low_density_pages":         low_density_pages,
        "pages_no_extracted_items":  pages_no_extracted_items,
        # 表格品質
        "empty_cell_rate":           empty_cell_rate,
        # 圖片
        "figures_total":             figures_total,
        "figures_materialized":      figures_materialized,
        # 綜合
        "estimated_info_loss_rate":  estimated,
        "qc_level":                  qc_level,
        "warnings":                  [],
        "errors":                    [],
    }


# ── 格式 Fallback 輔助函式 ────────────────────────────────────────────────────

def _spreadsheet_to_markdown(path: Path) -> str:
    """XLS / XLSX 轉 Markdown 表格（跳過圖表，只讀儲存格值）。"""
    suffix = path.suffix.lower()
    md_parts: list[str] = []

    if suffix == ".xls":
        import xlrd
        wb = xlrd.open_workbook(str(path))
        for sheet in wb.sheets():
            md_parts.append(f"## {sheet.name}")
            for row_idx in range(sheet.nrows):
                row = [str(sheet.cell_value(row_idx, c) or "") for c in range(sheet.ncols)]
                if any(v.strip() for v in row):
                    md_parts.append("| " + " | ".join(row) + " |")
    else:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            md_parts.append(f"## {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                vals = [str(v) if v is not None else "" for v in row]
                if any(v.strip() for v in vals):
                    md_parts.append("| " + " | ".join(vals) + " |")

    return "\n\n".join(md_parts)


# ── 主入口 ───────────────────────────────────────────────────────────────────

def convert_pdf_docling(
    pdf_path: Path,
    llm=None,
    category: str = "",
    output_dir: Path | None = None,
    describe_visuals: bool = False,
) -> dict:
    """Docling path：PDF → Docling 原始結構（schema-v3.0）。

    export_to_dict() 完整保留 Docling 輸出：
    - 每個 item 的精確 bbox
    - 表格的 cell-level 結構（row/col/span/header/bbox per cell）
    - formatting、hyperlink、charspan
    - 圖片用 bbox 精確裁切（比整頁更準確）

    Args:
        pdf_path       : PDF 路徑
        llm            : LLM 實例（有值時對 markdown 進行關鍵字萃取）
        category       : 文件類別
        output_dir     : 圖片輸出目錄（有值時用 bbox 精確裁切）
        describe_visuals: 保留供 extractor 介面一致性，不使用
    """
    from docling.document_converter import DocumentConverter
    from metadata_builder import build_metadata

    suffix = pdf_path.suffix.lower()

    # .doc → .docx conversion via LibreOffice
    if suffix == ".doc":
        import subprocess, tempfile
        with tempfile.TemporaryDirectory() as _tmpdir:
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "docx",
                 "--outdir", _tmpdir, str(pdf_path)],
                capture_output=True, timeout=60,
            )
            if result.returncode != 0:
                raise RuntimeError(
                    f"libreoffice conversion failed: {result.stderr.decode()}"
                )
            _converted = Path(_tmpdir) / (Path(pdf_path).stem + ".docx")
            if not _converted.exists():
                raise FileNotFoundError(f"Converted file not found: {_converted}")
            # Recurse with the converted .docx (won't enter this branch again)
            return convert_pdf_docling(
                _converted,
                llm=llm,
                category=category,
                output_dir=output_dir,
                describe_visuals=describe_visuals,
            )

    # XLS：docling 不支援，直接用 xlrd fallback
    if suffix == ".xls":
        try:
            markdown = _spreadsheet_to_markdown(pdf_path)
            raw = {}
        except Exception as exc:
            markdown = ""
            raw = {}

    else:
        # 1. Docling 轉換
        try:
            conv     = DocumentConverter()
            docling_result = conv.convert(pdf_path)
            doc      = docling_result.document
            raw      = doc.export_to_dict()
            markdown = doc.export_to_markdown()

        except Exception as exc:
            # XLSX 嵌入圖表導致 crash：用 openpyxl 讀純儲存格值
            if suffix == ".xlsx":
                try:
                    markdown = _spreadsheet_to_markdown(pdf_path)
                    raw = {}
                except Exception:
                    markdown = ""
                    raw = {}
            else:
                # 其他格式失敗：回傳錯誤結構
                page_count = 0
                if suffix == ".pdf":
                    try:
                        import fitz as _fitz
                        _doc = _fitz.open(str(pdf_path))
                        page_count = _doc.page_count
                        _doc.close()
                    except Exception:
                        pass
                metadata = build_metadata(
                    pdf_path=pdf_path, category=category,
                    extractor="docling", page_count=page_count,
                    markdown="", llm=llm,
                )
                metadata["confidence"] = {
                    "source":    "docling_structure",
                    "available": False,
                    "note":      "Docling conversion failed; no word-level confidence",
                }
                metadata["qc"] = {
                    "estimated_info_loss_rate": 1.0,
                    "qc_level":  "danger",
                    "warnings":  [],
                    "errors":    [f"DOCLING_CONVERSION_FAILED: {type(exc).__name__}: {exc}"],
                }
                return {
                    "schema_version": "v3.0",
                    "metadata":   metadata,
                    "data":       {},
                    "page_count": page_count,
                }

    # 2. 匯出完整 Docling 原始結構

    page_count = len(raw.get("pages", {})) or 1
    page_sizes = {k: v.get("size", {}) for k, v in raw.get("pages", {}).items()}

    # 3. 圖片裁切（使用 Docling bbox）
    # PDF：fitz 精確裁切；DOCX：用 Docling 內嵌的 base64 data URI 直接存檔
    pictures_raw = raw.get("pictures", [])
    is_pdf = pdf_path.suffix.lower() == ".pdf"
    is_docx = suffix == ".docx"
    if is_pdf and output_dir and pictures_raw:
        pictures = _materialize_pictures(pdf_path, pictures_raw, output_dir, page_sizes)
    elif is_docx and output_dir and pictures_raw:
        pictures = _materialize_pictures_from_uri(pictures_raw, output_dir, pdf_path.stem)
    else:
        pictures = [dict(p, path=None, sha256=None, has_image=False) for p in pictures_raw]
    raw["pictures"] = pictures

    # 3b. 整頁圖片（有圖片的頁面 + 無文字的頁面）
    page_images: dict[int, dict] = {}
    if is_pdf and output_dir:
        try:
            import fitz as _fitz
            # 有圖片的頁
            pic_pages = {
                prov.get("page_no")
                for pic in pictures if pic.get("has_image")
                for prov in pic.get("prov", [])
                if prov.get("page_no")
            }
            # 無文字的頁（掃描頁）
            qc_tmp = _compute_qc(raw, page_count, pictures, True)
            no_item_pages = set(qc_tmp.get("pages_no_extracted_items", []))
            full_pages = pic_pages | no_item_pages

            if full_pages:
                img_dir = output_dir / "figures"
                img_dir.mkdir(parents=True, exist_ok=True)
                doc = _fitz.open(str(pdf_path))
                try:
                    for pn in sorted(full_pages):
                        idx = pn - 1
                        if idx >= len(doc):
                            continue
                        pixmap    = doc[idx].get_pixmap(dpi=150)
                        png_bytes = pixmap.tobytes("png")
                        sha256    = _hashlib.sha256(png_bytes).hexdigest()
                        fpath     = img_dir / f"{pdf_path.stem}_p{pn}_full.png"
                        if not fpath.exists():
                            fpath.write_bytes(png_bytes)
                        page_images[pn] = {
                            "path":      str(fpath.relative_to(output_dir)),
                            "sha256":    sha256,
                            "has_image": True,
                        }
                finally:
                    doc.close()
        except Exception:
            pass  # 整頁圖片失敗不中斷主流程

    # 4. metadata
    title_text = next(
        (t.get("text", "") for t in raw.get("texts", []) if t.get("label") == "title"),
        None,
    )
    metadata = build_metadata(
        pdf_path=pdf_path, category=category,
        extractor="docling", page_count=page_count, title=title_text,
        markdown=markdown,
        llm=llm,
    )

    # 5. 組合 schema-v3.0 output（metadata + data 兩層）
    confidence = {
        "source":    "docling_structure",
        "available": False,
        "note":      "Layout/text/table extracted by Docling; no word-level confidence exposed",
    }
    qc = _compute_qc(raw, page_count, pictures, output_dir is not None)

    # 掃描型 PDF 偵測：所有頁面都沒有萃取出文字，且是 PDF
    if is_pdf and qc.get("pages_no_extracted_items"):
        no_item_count = len(qc["pages_no_extracted_items"])
        if no_item_count == page_count:
            qc["warnings"].append(
                "SCANNED_PDF_ALL_PAGES_EMPTY: docling 無法從掃描型 PDF 萃取文字。"
                "建議改用 azure_cu（支援 OCR）或 llm extractor。"
            )
        elif no_item_count > page_count // 2:
            qc["warnings"].append(
                f"SCANNED_PDF_PARTIAL: {no_item_count}/{page_count} 頁無文字萃取。"
                "部分頁面可能是掃描影像，建議改用 azure_cu。"
            )

    # 小字體偵測（< 8pt）：臨床關鍵數值可能精準度下降
    if is_pdf:
        try:
            import fitz as _fitz_sf
            _doc_sf = _fitz_sf.open(str(pdf_path))
            small_font_pages = []
            for _pi in range(_doc_sf.page_count):
                _rawdict = _doc_sf[_pi].get_text("rawdict")
                for _blk in _rawdict.get("blocks", []):
                    for _ln in _blk.get("lines", []):
                        for _sp in _ln.get("spans", []):
                            if _sp.get("size", 99) < 8 and _sp.get("text", "").strip():
                                small_font_pages.append(_pi + 1)
                                break
                        else:
                            continue
                        break
            _doc_sf.close()
            if small_font_pages:
                qc["warnings"].append(
                    f"SMALL_FONT_PAGES {sorted(set(small_font_pages))}: "
                    "部分頁面含 <8pt 文字（臨床分期表、劑量值等），OCR 精準度可能下降。"
                )
        except Exception:
            pass

    # Gap 3：JPG/PNG OCR 近乎空白 → 儲存原始圖作為 page_images fallback
    if suffix in {".jpg", ".jpeg", ".png"} and output_dir:
        meaningful = markdown.replace("<!-- image -->", "").replace("\n", "").strip()
        if len(meaningful) < 50:
            try:
                img_dir = output_dir / "figures"
                img_dir.mkdir(parents=True, exist_ok=True)
                dest = img_dir / f"{pdf_path.stem}_p1_full{suffix}"
                import shutil as _shutil
                _shutil.copy2(str(pdf_path), str(dest))
                sha256 = _hashlib.sha256(pdf_path.read_bytes()).hexdigest()
                page_images[1] = {
                    "path":      str(dest.relative_to(output_dir)),
                    "sha256":    sha256,
                    "has_image": True,
                }
                qc["warnings"].append("IMAGE_OCR_INSUFFICIENT_STORED_AS_PAGE_IMAGE")
            except Exception:
                pass

    # Gap 4：PPTX 空白輸出 → 用 python-pptx 萃取投影片圖片
    if suffix == ".pptx" and output_dir:
        meaningful = markdown.replace("<!-- image -->", "").replace("\n", "").strip()
        if len(meaningful) < 100:
            try:
                from pptx import Presentation as _Presentation
                prs = _Presentation(str(pdf_path))
                img_dir = output_dir / "figures"
                img_dir.mkdir(parents=True, exist_ok=True)
                for slide_idx, slide in enumerate(prs.slides, 1):
                    for shape in slide.shapes:
                        if shape.shape_type == 13:  # PICTURE
                            img_data = shape.image.blob
                            ext      = shape.image.ext or "png"
                            sha256   = _hashlib.sha256(img_data).hexdigest()
                            fname    = f"{pdf_path.stem}_p{slide_idx}_full.{ext}"
                            fpath    = img_dir / fname
                            fpath.write_bytes(img_data)
                            page_images[slide_idx] = {
                                "path":      str(fpath.relative_to(output_dir)),
                                "sha256":    sha256,
                                "has_image": True,
                            }
                            break  # 每張投影片取第一張圖
                qc["warnings"].append("PPTX_TEXT_EMPTY_SLIDES_STORED_AS_PAGE_IMAGES")
            except Exception:
                pass

    metadata["confidence"] = confidence
    metadata["qc"]         = qc

    # 取得 Docling 套件版本（安裝後可讀取；讀取失敗則記為 ">=2.0"）
    try:
        from importlib.metadata import version as _pkg_version
        _docling_version = _pkg_version("docling")
    except Exception:
        _docling_version = ">=2.0"

    return {
        "schema_version": "v3.0",
        "metadata": metadata,
        "data": {
            "markdown":        markdown,
            "texts":           raw.get("texts",           []),
            "tables":          raw.get("tables",          []),
            "pictures":        pictures,
            "pages":           raw.get("pages",           {}),
            "groups":          raw.get("groups",          []),
            "key_value_items": raw.get("key_value_items", []),
            "body":            raw.get("body",            {}),
            "page_images":     page_images,
            "furniture":       raw.get("furniture",       {}),
        },
        "page_count": page_count,
        "extractor_metadata": {
            "tool":                "docling",
            "analyzer_id":         None,
            "api_version":         None,
            "parsing_mode":        "xml_native",
            "version":             _docling_version,
            "confidence_source":   None,
            "per_cell_confidence": False,
            "confidence_note": (
                "Docling 從 XML 直接解析，不經 OCR，因此無信心分數；"
                "結構正確性由 empty_cell_rate 與 garbled_char_rate 間接評估"
            ),
            "header_flags": (
                "column_header / row_header 布林欄位由 Docling 從 XML 解析，非 heuristic"
            ),
            "bounding_box_format": (
                "Docling BOTTOMLEFT bbox（l/t/r/b，pt 單位）"
            ),
            "known_limitation": (
                "複雜三層以上 merged header 的 row_span 計算偶有已知 Docling GitHub issues；"
                "部分 item（embedded object / floating element）的 prov[] 陣列可能為空或缺少 page_no 欄位，"
                "下游應以繼承前一 item 頁碼或標記 page_unknown 作為 fallback"
            ),
            "fallback_reason":     None,
        },
    }
